import shutil
from pptx import Presentation

source_file = 'templates/Consulting_Template_2.pptx'
destination_file = 'outputs/Copied_File_4.pptx'

shutil.copy(source_file, destination_file)

presentation = Presentation('outputs/Copied_File_4.pptx')

slide_count = len(presentation.slides)
print(f"Number of slides: {slide_count}")

layout_names = []

for slide_master in presentation.slide_masters:
    for slide_layout in slide_master.slide_layouts:
        layout_names.append(slide_layout.name)

# for layout_name in layout_names:
#     print(layout_name)

print("Select a slide layout by entering its index:")
for i, layout_name in enumerate(layout_names):
    print(f"{i + 1}. {layout_name}")

selected_layout_index = int(input("Enter the index of the slide layout you want to add: ")) - 1
selected_layout = presentation.slide_layouts[selected_layout_index]
print(f"Selected layout: {selected_layout}")

num_slides_to_add = int(input("How many slides would you like to add? "))

for _ in range(num_slides_to_add):
    new_slide = presentation.slides.add_slide(selected_layout)
    print(f"New slide: {new_slide}")

slide_count = len(presentation.slides)
print(f"Number of slides: {slide_count}")

presentation.save('outputs/Copied_File_4.pptx')
