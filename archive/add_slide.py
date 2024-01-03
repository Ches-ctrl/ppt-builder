from pptx import Presentation

prs = Presentation('test.pptx')

title_slide_layout = prs.slide_layouts[0]
new_slide = prs.slides.add_slide(title_slide_layout)

slide_count = len(prs.slides)
print(f"Number of slides: {slide_count}")

prs.save('test.pptx')
