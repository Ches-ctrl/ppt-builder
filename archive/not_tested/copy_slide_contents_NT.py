from pptx import Presentation

# Load the source presentation
source_presentation = Presentation('templates/Consulting_Template.pptx')

# Create a new presentation to copy the slide into
destination_presentation = Presentation()

# Specify the source slide index (0-based) to copy from
source_slide_index = 0

# Specify the destination slide layout (0 for the default layout)
destination_slide_layout_index = 0

# Copy the source slide to the destination presentation
source_slide = source_presentation.slides[source_slide_index]
destination_slide_layout = destination_presentation.slide_layouts[destination_slide_layout_index]
destination_slide = destination_presentation.slides.add_slide(destination_slide_layout)

# Copy all elements from the source slide to the destination slide
for source_element in source_slide.shapes:
    # Clone the source element and add it to the destination slide
    destination_element = source_element.clone()
    destination_slide.shapes.add_element(destination_element)

# Save the destination presentation
destination_presentation.save('outputs/Copied_Slide_Contents_1.pptx')
