from pptx import Presentation

presentation = Presentation('templates/Delete_Slide_1.pptx')

presentation_length = len(presentation.slides)
print(f"Presentation length: {presentation_length}")

user_input = int(input("Enter the index of the slide you want to delete: "))
slide_index_to_delete = user_input - 1
slide_id_to_delete = presentation.slides[slide_index_to_delete].slide_id
print(f"Slide ID to delete: {slide_id_to_delete}")

sldId_lst = presentation.slides._sldIdLst

for sldId in sldId_lst:
    if sldId.id == slide_id_to_delete:
        sldId_lst.remove(sldId)
        print(f"Slide ID {slide_id_to_delete} removed from the slide ID list.")
        break

presentation_length = len(presentation.slides)
print(f"Presentation length: {presentation_length}")

presentation.save('outputs/modified_presentation.pptx')
