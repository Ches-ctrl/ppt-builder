from pptx import Presentation

new_presentation = Presentation()
template_files = ['templates/template_1.pptx', 'templates/template_2.pptx', 'templates/template_3.pptx']

for template_file in template_files:
    template = Presentation(template_file)

    for slide in template.slides:
        slide_layout = new_presentation.slide_layouts.get_by_name(slide.slide_layout.name)
        new_slide = new_presentation.slides.add_slide(slide_layout)

new_presentation.save('outputs/combined_presentation.pptx')
