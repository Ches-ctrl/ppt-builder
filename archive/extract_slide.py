from pptx import Presentation

source_prs = Presentation('test.pptx')
destination_prs = Presentation()

first_slide_layout = source_prs.slides[0].slide_layout

destination_prs.slides.add_slide(first_slide_layout)

slide_count = len(destination_prs.slides)
print(f"Number of slides: {slide_count}")

destination_prs.save('output.pptx')
