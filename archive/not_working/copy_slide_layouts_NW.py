from pptx import Presentation

source_presentation = Presentation('templates/Consulting_Template.pptx')
destination_presentation = Presentation('outputs/destination_presentation.pptx')

for source_slide_master in source_presentation.slide_masters:
    for source_slide_layout in source_slide_master.slide_layouts:
        print(f"Source slide layout: {source_slide_layout}")
        existing_layout = destination_presentation.slide_master.slide_layouts.get_by_name(source_slide_layout.name)
        print(f"Existing layout: {existing_layout}")

        if existing_layout:
            destination_presentation.slide_master.slide_layouts.remove(existing_layout)

        new_slide_layout = destination_presentation.slide_master.slide_layouts.add_slide_layout(source_slide_layout.name)

destination_presentation.save('modified_destination_presentation.pptx')
