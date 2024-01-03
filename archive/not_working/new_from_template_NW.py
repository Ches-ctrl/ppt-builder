from pptx import Presentation

new_presentation = Presentation()
template_presentation = Presentation('templates/EYP_Template.pptx')

for slide in template_presentation.slides:
    slide_layout_name = slide.slide_layout.name
    print(f"Slide layout name: {slide_layout_name}")
    if new_presentation.slide_layouts.get_by_name(slide_layout_name):
        slide_layout = new_presentation.slide_layouts.get_by_name(slide_layout_name)
        new_slide = new_presentation.slides.add_slide(slide_layout)
    else:
        print(f"Warning: Slide layout '{slide_layout_name}' not found in the new presentation.")

new_presentation.save('outputs/output_test_1.pptx')
