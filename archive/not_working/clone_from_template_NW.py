from pptx import Presentation

new_presentation = Presentation()
template_presentation = Presentation('templates/Consulting_Template.pptx')

for template_slide in template_presentation.slides:
    new_slide = new_presentation.slides.add_slide(template_slide.slide_layout)

new_presentation.save('outputs/output_test_1.pptx')
