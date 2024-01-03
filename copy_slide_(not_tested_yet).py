from pptx import Presentation

# Load the source and destination presentations
source_presentation = Presentation('source_presentation.pptx')
destination_presentation = Presentation('destination_presentation.pptx')

# Choose the slide to copy (e.g., slide 1)
source_slide = source_presentation.slides[0]

# Create a new slide in the destination presentation
new_slide_layout = destination_presentation.slide_layouts[source_slide.slide_layout.name]
new_slide = destination_presentation.slides.add_slide(new_slide_layout)

# Copy content from the source slide to the new slide
for source_shape in source_slide.shapes:
    if source_shape.has_text_frame:
        new_shape = new_slide.shapes.add_shape(
            source_shape.auto_shape_type,
            source_shape.left, source_shape.top,
            source_shape.width, source_shape.height
        )
        new_shape.text = source_shape.text

# Save the destination presentation
destination_presentation.save('output_presentation.pptx')
